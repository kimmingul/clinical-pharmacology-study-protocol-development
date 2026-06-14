#!/usr/bin/env python3
"""Generate the 20-slide PPTX for the talk
"생성형 AI를 활용한 임상개발전략의 효율화: Protocol 작성 및 DDI 설계 사례".

Content + full speaker scripts come from presentation_genAI_protocol_DDI_ko.md.
The on-slide text is intentionally sparse; the full word-for-word script goes
into each slide's NOTES pane. Theme: navy + grey, red accent for "trap" slides.

Run:  python build_presentation_pptx.py
Out:  presentation_genAI_protocol_DDI_ko.pptx  (editable in PowerPoint)
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette --------------------------------------------------------------
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
NAVY2 = RGBColor(0x2E, 0x5A, 0x88)
GREY = RGBColor(0x5A, 0x5A, 0x5A)
LIGHT = RGBColor(0xEE, 0xF1, 0xF6)
RED = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x26, 0x2B)
FONT = "맑은 고딕"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def _set_font(run, size, color=DARK, bold=False, font=FONT, italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font)


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def _rect(slide, l, t, w, h, fill, line=None, rounded=False):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shp, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def _para(tf, text, size, color=DARK, bold=False, bullet=False, space=8,
          align=PP_ALIGN.LEFT, italic=False, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    r = p.add_run()
    r.text = ("•  " + text) if bullet else text
    _set_font(r, size, color, bold, italic=italic)
    return p


prs = Presentation()
prs.slide_width, prs.slide_height = EMU_W, EMU_H
BLANK = prs.slide_layouts[6]


def base(num, section, title, accent=False):
    s = prs.slides.add_slide(BLANK)
    _rect(s, 0, 0, 13.333, 1.2, NAVY)
    if accent:
        _rect(s, 0, 1.2, 13.333, 0.09, RED)
    tf = _box(s, 0.55, 0.18, 12.2, 0.95)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _para(tf, title, 26, WHITE, bold=True, first=True, space=0)
    ftf = _box(s, 0.55, 7.0, 12.2, 0.4)
    _para(ftf, f"{section}    ·    {num} / 20", 10, GREY, first=True, space=0)
    return s


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------------------
# slide data: each entry = dict(num, sec, title, accent, render(fn), notes)
# ---------------------------------------------------------------------------
CONTENT_TOP = 1.6


def bullets(items, top=CONTENT_TOP, size=18, left=0.9, width=11.5):
    def f(s):
        tf = _box(s, left, top, width, 5.0)
        for i, it in enumerate(items):
            if isinstance(it, tuple):  # (text, indent)
                txt, lvl = it
                _para(tf, txt, size - (2 * lvl), DARK if lvl == 0 else GREY,
                      bold=(lvl == 0), bullet=(lvl == 0), first=(i == 0), space=10)
            else:
                _para(tf, it, size, DARK, bold=True, bullet=True,
                      first=(i == 0), space=12)
    return f


def big(text, sub=None, color=NAVY):
    def f(s):
        tf = _box(s, 1.0, 2.6, 11.3, 2.2)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _para(tf, text, 34, color, bold=True, first=True, align=PP_ALIGN.CENTER, space=14)
        if sub:
            _para(tf, sub, 18, GREY, align=PP_ALIGN.CENTER, space=0)
    return f


def kpi(cards, caption=None):
    def f(s):
        n = len(cards)
        gap = 0.5
        w = (12.3 - gap * (n - 1)) / n
        x = 0.5
        for big_t, small_t, col in cards:
            _rect(s, x, 2.0, w, 2.4, LIGHT, rounded=True)
            tf = _box(s, x + 0.15, 2.2, w - 0.3, 2.0)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            _para(tf, big_t, 30, col, bold=True, first=True, align=PP_ALIGN.CENTER, space=8)
            _para(tf, small_t, 15, GREY, align=PP_ALIGN.CENTER, space=0)
            x += w + gap
        if caption:
            ctf = _box(s, 0.5, 4.7, 12.3, 0.8)
            _para(ctf, caption, 17, RED, bold=True, first=True, align=PP_ALIGN.CENTER, space=0)
    return f


def pillars(cards):
    def f(s):
        n = len(cards)
        gap = 0.4
        w = (12.3 - gap * (n - 1)) / n
        x = 0.5
        for i, (head, desc) in enumerate(cards):
            _rect(s, x, 2.0, w, 3.2, NAVY if i % 2 == 0 else NAVY2, rounded=True)
            tf = _box(s, x + 0.2, 2.25, w - 0.4, 2.8)
            _para(tf, head, 17, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER, space=10)
            _para(tf, desc, 13, RGBColor(0xDD, 0xE4, 0xEE), align=PP_ALIGN.CENTER, space=0)
            x += w + gap
        ctf = _box(s, 0.5, 5.45, 12.3, 0.7)
        _para(ctf, "그리고 이 모든 것의 가운데에 — 사람의 승인 게이트", 16, GREY,
              bold=True, first=True, align=PP_ALIGN.CENTER, space=0)
    return f


def two_col(left_title, left_items, right_title, right_items,
            lcolor=GREY, rcolor=NAVY):
    def f(s):
        _rect(s, 0.5, 1.7, 6.0, 4.9, RGBColor(0xF6, 0xF1, 0xF0), rounded=True)
        _rect(s, 6.9, 1.7, 6.0, 4.9, LIGHT, rounded=True)
        lt = _box(s, 0.7, 1.85, 5.6, 0.6)
        _para(lt, left_title, 18, lcolor, bold=True, first=True, space=0)
        ltf = _box(s, 0.7, 2.5, 5.6, 3.9)
        for i, it in enumerate(left_items):
            _para(ltf, it, 14, DARK, bullet=True, first=(i == 0), space=9)
        rt = _box(s, 7.1, 1.85, 5.6, 0.6)
        _para(rt, right_title, 18, rcolor, bold=True, first=True, space=0)
        rtf = _box(s, 7.1, 2.5, 5.6, 3.9)
        for i, it in enumerate(right_items):
            _para(rtf, it, 14, DARK, bullet=True, first=(i == 0), space=9)
    return f


def table17():
    headers = ["설계 단계", "Before: 범용 AI", "After: harness 기반"]
    rows = [
        ["① 질문 정의", "곧바로 'DDI 계획서' 작성 시작",
         "먼저 해체: 단방향인가 양방향인가? 부터 결정"],
        ["② 기전 전제", "'상호작용 가능' 일반 서술",
         "clarithromycin=CYP3A4 저해 → tegoprazan↑ / tegoprazan=위산억제 → clarithromycin·14-OH↑ ⇒ 양방향"],
        ["③ 설계 선택", "open-label, fixed-sequence 막연",
         "양방향 평가 → 6-sequence, 3-period crossover (carry-over 균형)"],
        ["④ 평가변수", "Cmax, AUC 나열",
         "양방향 각각 GMR·90% CI. 실제: 병용 시 tegoprazan AUC ~2.5배↑"],
        ["⑤ 표본수", "숫자 하나(가정·검증 없음)",
         "CV·GMR·검정력 명시 + 양방향 IUT + 검증된 계산 코드"],
        ["⑥ 안전성·채혈", "일반 나열",
         "clarithromycin 관련(간·QT·병용금기) + Tmax·소실기 반영"],
        ["⑦ 검토·근거", "없음",
         "다중 관점 검토 + 출처(PMID 37440779) + 사람 승인 게이트"],
    ]

    def f(s):
        rows_n = len(rows) + 1
        tbl = s.shapes.add_table(rows_n, 3, Inches(0.45), Inches(1.55),
                                 Inches(12.45), Inches(5.3)).table
        tbl.columns[0].width = Inches(2.0)
        tbl.columns[1].width = Inches(4.1)
        tbl.columns[2].width = Inches(6.35)
        for c, h in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = h
            _set_font(r, 13, WHITE, bold=True)
        for ri, row in enumerate(rows, start=1):
            for ci, val in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = cell.text_frame.paragraphs[0]
                cell.text_frame.word_wrap = True
                r = p.add_run()
                r.text = val
                col = GREY if ci == 1 else (NAVY if ci == 2 else DARK)
                _set_font(r, 10.5, col, bold=(ci == 0))
    return f


def timeline(steps):
    def f(s):
        n = len(steps)
        gap = 0.3
        w = (12.3 - gap * (n - 1)) / n
        x = 0.5
        for i, (yr, txt, hot) in enumerate(steps):
            _rect(s, x, 2.6, w, 1.7, RED if hot else NAVY, rounded=True)
            tf = _box(s, x + 0.1, 2.7, w - 0.2, 1.5)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            _para(tf, yr, 16, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER, space=6)
            _para(tf, txt, 11, RGBColor(0xE8, 0xEC, 0xF2), align=PP_ALIGN.CENTER, space=0)
            x += w + gap
        ctf = _box(s, 0.5, 4.7, 12.3, 0.9)
        _para(ctf, "규제가 묻는 건 '모델 성능'이 아니라 '어떤 체계 속에서 쓰는가'", 18,
              NAVY, bold=True, first=True, align=PP_ALIGN.CENTER, space=0)
    return f


# ---------------------------------------------------------------------------
SLIDES = [
    dict(sec="도입", title="생성형 AI를 활용한 임상개발전략의 효율화", accent=False,
         render=big("Protocol 작성 및 DDI 설계 사례",
                    "AI를 '신뢰할 수 있는 동료'로 만드는 법  ·  한국임상개발협회", WHITE if False else NAVY),
         notes="안녕하십니까. 오늘은 생성형 AI를 임상개발에 어떻게 활용할 것인가를 말씀드립니다. 다만 결론부터 말씀드리면, 오늘의 핵심은 'AI로 무엇을 자동화하느냐'가 아니라 '그 결과를 어떻게 믿고 쓰느냐'입니다. 임상개발 전문가의 시선에서 30분간 함께 보시겠습니다."),
    dict(sec="도입", title="질문 하나로 시작합니다", accent=False,
         render=big("신입 연구원에게 첫 출근날\n'DDI 계획서 써오세요' 하시나요?",
                    "그런데 우리는 AI에게는 그렇게 시키고 있습니다."),
         notes="질문 하나로 시작하겠습니다. 신입 연구원이 첫 출근한 날 곧바로 '약물상호작용 시험 계획서 작성해 오세요' 하실 분 계신가요? 아마 없으실 겁니다. 그 신입이 약동학을 몰라서가 아닙니다. 우리 조직의 템플릿, 규제 표현, 시험 유형별 필수 항목, 검토 문화를 아직 모르기 때문입니다. 지식이 아니라 맥락이 없어서죠. 그런데 우리는 AI에게는 정확히 그렇게 시킵니다. 오늘은 그 간극과, 그 간극을 메우는 법을 말씀드립니다."),
    dict(sec="도입", title="오늘의 3가지 메시지", accent=False,
         render=bullets([
             "① AI는 왜 '혼자선' 부족한가",
             "② '신뢰'는 어떻게 만드는가",
             "③ '우리 조직'은 무엇부터 할까",
         ], top=2.2, size=24),
         notes="오늘 딱 세 가지만 가져가시면 됩니다. 첫째, AI는 왜 혼자서는 부족한가. 둘째, 그렇다면 신뢰는 어떻게 만드는가. 셋째, 우리 조직은 내일 무엇부터 시작하면 되는가. 이 세 가지를 DDI 설계라는 구체적 사례로 풀어가겠습니다."),
    dict(sec="Act 1 · 기대와 함정", title="생성형 AI, 1분 이해", accent=False,
         render=bullets([
             "생성형 AI = '방대한 글을 읽고 다음에 올 가장 그럴듯한 말을 잇는' 모델",
             ("핵심 성질: 사실(fact)이 아니라 '그럴듯함'을 최적화한다", 1),
             "프롬프트 = AI에게 주는 지시문",
             "환각 = 그럴듯하지만 사실이 아닌 답을 자신 있게 만드는 현상",
             "에이전트 = 특정 역할을 맡은 AI 한 명",
         ], size=18),
         notes="본론 전에 1분만 개념을 정리하겠습니다. 생성형 AI는 본질적으로 '방대한 글을 학습해 다음에 올 가장 그럴듯한 말을 잇는' 기계입니다. 중요한 단어는 '그럴듯한'입니다. AI는 사실이 아니라 그럴듯함을 최적화합니다. 세 단어만 기억해 주십시오. 프롬프트는 지시문, 환각은 그럴듯한 거짓, 에이전트는 역할을 맡은 AI 한 명입니다."),
    dict(sec="Act 1 · 기대와 함정", title="데모: 범용 AI에게 DDI 계획서를 시키면", accent=False,
         render=bullets([
             "프롬프트: 'Tegoprazan–Clarithromycin DDI 계획서를 작성하라'",
             ("→ 30초 만에 목적·설계·Cmax/AUC·안전성이 '그럴듯하게' 채워진 초안", 1),
             "[이 자리에 실제 화면 캡처를 넣으세요]",
         ], size=18),
         notes="실제로 해봤습니다. 범용 AI에게 'Tegoprazan과 Clarithromycin DDI 계획서를 작성하라'고 시키면 30초도 안 되어 나옵니다. 목적, 설계, Cmax와 AUC, 안전성까지 채워져 있습니다. 빠르고 매끄럽죠. 많은 분이 여기서 'AI 되네' 하십니다. 그런데 우리는 전문가입니다. 한 겹만 벗겨 보겠습니다."),
    dict(sec="Act 1 · 기대와 함정", title="전문가의 눈: 빠진 '설계 논리'", accent=True,
         render=bullets([
             "왜 fixed-sequence인가?",
             "clarithromycin 투여 기간의 근거는?",
             "채혈 시점은 무엇을 기준으로?",
             "단방향인가, 양방향인가?",
             "표본수 가정은? washout 정당화는?",
         ], size=20, left=1.2),
         notes="전문가의 질문을 던져보겠습니다. 왜 fixed-sequence죠? clarithromycin은 며칠을 투여해야 정상상태 저해에 도달하나요? 채혈 시점 근거는? 그리고 가장 중요한 — 단방향입니까 양방향입니까? 표본수는 어떤 가정으로 계산했죠? 형식은 다 갖췄는데 '왜 그렇게 설계했는지'가 비어 있습니다. 임상개발 문서의 본질이 바로 그 '왜'인데 말입니다."),
    dict(sec="Act 1 · 기대와 함정", title="핵심 함정: 못 쓰는 게 아니라, 그럴듯하게 틀린다", accent=True,
         render=kpi([
             ("52 → 28", "AI가 만든 표본수 vs 정답 (2배 과대)", RED),
             ("PMID ❓", "존재하지 않는 환각 인용", RED),
         ], caption="이대로 IND 신청서에 들어갔다면?"),
         notes="오늘 첫 번째 핵심 슬라이드입니다. AI의 진짜 위험은 못 하는 게 아니라 자신 있게 틀리는 것입니다. 두 사례를 보여드립니다. 첫째, 제가 이 시스템을 만들며 직접 겪은 일입니다. AI의 도움으로 짠 표본수 코드가 어떤 조건에서 52명이 필요하다 했는데, 정확히 다시 계산하니 28명이 정답이었습니다. 두 배 과대 산출이죠. 둘째, 그럴듯한 PMID를 제시했는데 존재하지 않는 번호였습니다 — 환각입니다. 이게 검토 없이 IND에 들어갔다면? 표본수가 두 배면 비용·시간, 그리고 불필요하게 더 많은 자원자가 위험에 노출됩니다. 그리고 이런 오류는 '문장이 유창한가'로는 절대 안 보입니다. 오직 검증으로만 드러납니다."),
    dict(sec="Act 1 · 기대와 함정", title="왜 틀리나: protocol은 '판단의 집합체'", accent=False,
         render=two_col("AI가 본 것", ["하나의 '문서'", "그럴듯한 문장의 나열"],
                        "실제 protocol",
                        ["약물 기전", "설계", "평가변수", "표본수", "안전성", "규제 판단의 응축"]),
         notes="왜 이런 일이 생길까요? AI는 계획서를 '하나의 문서'로 봅니다. 그런데 계획서는 문서가 아니라 판단의 집합체입니다. 누구를 넣고 뺄지, 어떤 설계를 쓸지, 무엇을 1차 변수로 둘지, 몇 명을 모집할지 — 수많은 판단이 한 문서에서 동시에 해결됩니다. 신입에게 첫날 계획서를 안 맡기는 이유와 똑같죠. 지식이 아니라 맥락이 없어서. 그렇다면 답도 분명합니다 — 신입을 교육하듯 AI에게 맥락을 갖춰주면 됩니다."),
    dict(sec="Act 2 · 신뢰 구조", title="해법: AI를 위한 '실행환경(harness)'", accent=False,
         render=two_col("신입 교육", ["조직 템플릿", "작성 기준", "선배의 검토", "계산 도구"],
                        "harness(실행환경)",
                        ["프로토콜 템플릿", "규제 라이브러리", "검토 단계", "검증된 계산 코드"]),
         notes="그 '맥락을 갖춘 구조'를 harness, 우리말로 실행환경이라 부르겠습니다. 핵심은 '좋은 모델보다 좋은 사용 구조'입니다. 신입을 교육할 때 무엇을 주나요? 템플릿, 기준, 선배 검토, 도구. harness도 똑같습니다. AI에게 템플릿, 규제 라이브러리, 검토 단계, 검증된 계산 코드를 갖춰줍니다. harness는 부가기능이 아니라 AI를 작동 가능하게 만드는 본체입니다."),
    dict(sec="Act 2 · 신뢰 구조", title="신뢰의 4기둥", accent=False,
         render=pillars([
             ("① 역할 분업", "실제 팀처럼 전공별로 나눔"),
             ("② 근거 추적", "모든 주장에 출처 의무"),
             ("③ 계산 검증", "코드로, 그리고 코드도 검증"),
             ("④ 검토·기록", "다중 검토 + provenance"),
         ]),
         notes="신뢰는 막연한 구호가 아니라 네 기둥으로 구조화됩니다. 역할 분업, 근거 출처, 계산 검증, 다중 검토와 기록. 그리고 가운데에 사람의 승인 게이트가 있습니다. 하나씩 보겠습니다."),
    dict(sec="Act 2 · 신뢰 구조", title="기둥①: 역할 기반 분업 (가상 임상개발팀)", accent=False,
         render=bullets([
             "임상약리 · 중개의학 · 규제 · 임상의 · 통계 · 작성 · 동의서 · QA = 8개 역할",
             "한 AI에 전부 맡기지 않고 '가상 CRO팀'처럼 분업",
             ("특히 '작성하는 AI'와 '검토하는 AI'를 분리 — 자기 글을 자기가 통과시키지 않게", 0),
         ], size=18),
         notes="첫 번째 기둥, 역할 분업입니다. 한 AI에 전부 시키지 않고 실제 팀처럼 나눕니다. 약동학의 임상약리, 바이오마커의 중개의학, 규제, 안전성의 임상의, 표본수의 통계, 그리고 작성과 검토가 따로 있습니다. 특히 작성하는 AI와 검토하는 AI를 분리한 게 중요합니다. 초안을 만드는 사고와 오류를 찾는 사고는 전혀 다른 작업이니까요. 자기 글을 자기가 통과시키지 않게 하는 겁니다."),
    dict(sec="Act 2 · 신뢰 구조", title="기둥②: 모든 주장에 '출처' 의무 (환각 차단)", accent=False,
         render=bullets([
             "PMID·NCT·가이드라인 없으면 → '[출처 미확인]' 표시, 날조 금지",
             ("PMID는 PubMed로, NCT는 ClinicalTrials.gov로 자동 검증 가능", 1),
             "자료원: PubMed · ClinicalTrials.gov · FDA 라벨 · MFDS 승인현황 · 가이드라인 라이브러리(ICH/FDA/EMA/MFDS)",
         ], size=17),
         notes="두 번째 기둥, 근거의 추적입니다. 환각 PMID를 막는 방법은 규칙을 강제하는 겁니다. 모든 주장에 PMID, NCT, 가이드라인 출처를 달게 하고, 확인 안 된 건 '출처 미확인'으로 표시, 없는 번호 생성은 금지. 그리고 PMID는 PubMed로, NCT는 ClinicalTrials.gov로 자동 검증까지 됩니다. 여기에 FDA 라벨, 식약처 승인현황, 가이드라인 라이브러리를 우선 참고하게 합니다. 임상개발에선 그럴듯한 문장보다 검증 가능한 근거가 중요하니까요."),
    dict(sec="Act 2 · 신뢰 구조", title="기둥③: 계산은 코드로, 그리고 코드도 검증한다", accent=True,
         render=kpi([
             ("52 → 28", "공식 오적용 (2배 과대)", RED),
             ("24 → 32", "GMR=1.0: 정규근사 24명은 실제 검정력 0.64", RED),
             ("✓ 검증", "PowerTOST·시뮬레이션·자동시험 73개", NAVY),
         ], caption="'코드로 옮겼다'가 아니라 '코드를 검증했다'에서 신뢰가 나온다"),
         notes="세 번째 기둥, 오늘 가장 중요한 슬라이드입니다. 표본수·초기용량 같은 숫자는 AI가 말로 내놓게 하면 안 되고, 검증된 계산 코드가 산출하게 해야 합니다. 그런데 한 걸음 더 — '코드로 옮겼다'가 끝이 아니라 그 코드 자체를 검증해야 합니다. 52 대 28이 바로 코드의 공식 오적용이었습니다. 더 미묘한 경우도 있었습니다. 기대 GMR이 정확히 1.0일 때 흔한 정규근사가 24명이면 된다 했는데, 실제 검정력은 0.64밖에 안 됐습니다. 정확한 방법으로는 32명이 필요했죠. 이런 건 눈으로 안 보입니다. 그래서 정확한 통계 방법으로 교체하고 표준 소프트웨어, 모의실험, 73개 자동시험으로 대조했습니다. 결론 — 검증되지 않은 계산은 사람이 짰든 AI가 짰든 똑같이 위험합니다."),
    dict(sec="Act 2 · 신뢰 구조", title="기둥④: 다중 검토 + 기록(provenance)", accent=False,
         render=bullets([
             "서로 다른 AI 3종이 '각자 다른' 결함을 발견 (통계 / 문서 정합성 / 배포·사용성)",
             ("→ 한 모델만 썼다면 못 잡았을 결함이 교차검토로 표면화", 1),
             "단계별 기록: 누가 · 어떤 모델 · 입력/출력 해시(SHA-256) · 시각 → 재현성·규제 방어",
             ("단, 이것도 자동 통과가 아니라 '사람 검토 전 결함 표면화'일 뿐", 1),
         ], size=17),
         notes="네 번째 기둥, 다중 검토와 기록입니다. 이 시스템을 점검할 때 서로 다른 세 AI에게 같은 산출물을 독립 검토시켰더니, 같은 결론을 반복한 게 아니라 서로 다른 결함을 찾았습니다. 한 모델은 통계 오류를, 다른 모델은 문서 경로 불일치를, 또 다른 모델은 배포·사용성 문제를요. 한 모델만 썼다면 못 잡았을 것들입니다. 그리고 모든 단계마다 누가 어떤 모델로 어떤 입력을 받아 무엇을 만들었는지를 해시와 시각까지 기록합니다. 재현성과 규제 방어를 위해서요. 다만 이 다중 검토도 사람 검토를 대체하지 않습니다 — 전처리일 뿐입니다."),
    dict(sec="Act 2 · 신뢰 구조", title="규제도 같은 곳을 가리킨다", accent=False,
         render=timeline([
             ("2023", "논의 출발", False),
             ("2024", "운영 원칙", False),
             ("2025", "신뢰성의 언어", False),
             ("2026.1", "FDA·EMA 공동원칙", True),
             ("2026.6", "ICH E6(R3) Annex 2 Step 4", True),
         ]),
         notes="이게 저 혼자 생각일까요? 아닙니다. 규제기관이 같은 곳을 가리킵니다. 2023년 논의 시작, 2024년 운영원칙, 2025년 '신뢰할 수 있는 AI 활용', 2026년 1월 FDA·EMA 공동원칙. 같은 시기 ICH E6 개정판도 진화했습니다. 강조하는 게 목적·맥락·위험기반·데이터품질·문서화, 그리고 인간 감독과 책임입니다. 이게 사실상 좋은 harness의 원칙과 똑같습니다. 규제는 막연히 금지·허용하지 않고 '어떤 체계 속에서 쓰는가'를 묻습니다."),
    dict(sec="Act 3 · DDI 사례", title="사례: Tegoprazan + Clarithromycin 양방향 DDI", accent=False,
         render=bullets([
             "실제 출판 연구 — PMID 37440779, NCT02052336 (Transl Clin Pharmacol 2023)",
             "설계: open-label, 6-sequence, 3-period crossover",
             "기전: clarithromycin(CYP3A4 저해) ↔ tegoprazan(위산 분비 억제) ⇒ 양방향",
             ("결과: 병용 시 tegoprazan AUCss,τ 약 2.5배↑ (Css,max 약 1.6배↑)", 1),
         ], size=18),
         notes="이제 구체 사례로 묶어보겠습니다. Tegoprazan과 Clarithromycin 상호작용입니다. 가상이 아니라 실제 출판된 연구입니다. 왜 좋은 교육 사례냐면 상호작용이 양방향이기 때문입니다. Clarithromycin은 CYP3A4 저해제라 그 효소로 대사되는 tegoprazan 노출을 올리고, 반대로 tegoprazan은 위산 분비를 억제해 clarithromycin과 활성 대사체의 위내 안정성·흡수에 영향을 줍니다. 실제 병용 시 tegoprazan 노출이 약 2.5배 증가했습니다. 이런 양방향에서는 설계 선택이 결과를 가릅니다."),
    dict(sec="Act 3 · DDI 사례", title="Before(범용 AI) vs After(harness) — 같은 질문, 다른 결과", accent=False,
         render=table17(),
         notes="이 표 한 장이 오늘 발표의 압축입니다. 왼쪽 범용 AI, 오른쪽 구조화된 접근. 질문 정의부터 다릅니다. 범용 AI는 바로 쓰고, 구조화된 접근은 먼저 단방향이냐 양방향이냐를 묻습니다. 이 한 질문이 다 바꿉니다. 기전에서 양방향이라는 결론에 도달하니 설계도 6-시퀀스 3-기간 교차설계가 됩니다. 실제 연구가 정확히 이걸 골랐죠. 평가변수도 양방향 각각의 GMR과 90% 신뢰구간, 표본수도 그 검증된 코드로 양방향 통계까지 반영합니다. 핵심은 오른쪽이 더 똑똑한 AI라서가 아니라 질문을 구조화했기 때문입니다."),
    dict(sec="Act 3 · DDI 사례", title="'효율화'의 실체: 무엇이 빨라지고, 무엇은 사람이 지키나", accent=False,
         render=two_col("AI가 가속 (주 → 일 단위)",
                        ["자료 조사", "유사시험 비교표", "초안 작성", "검토 전처리"],
                        "사람이 지킴 (자동화 금지)",
                        ["설계의 핵심 판단", "승인 게이트(자료·시놉시스)", "최종 검증"],
                        rcolor=RED),
         notes="그럼 '효율화'란 정확히 뭘까요. 솔직히 구분해 드리겠습니다. AI가 빠르게 해주는 건 자료 조사, 비교표, 초안, 검토 사전정리입니다 — 주 단위를 일 단위로 줄입니다. 그런데 사람이 반드시 지켜야 하는 것 — 설계 핵심 판단, 승인 게이트, 최종 검증 — 은 절대 자동화하지 않습니다. 효율화는 사람을 빼는 게 아니라, 반복을 줄여 사람이 판단에 집중하게 하는 것입니다."),
    dict(sec="마무리", title="우리 조직은 무엇부터? (실무 제언 5)", accent=False,
         render=bullets([
             "① 작은 use case부터 (배경조사 요약 · 비교표 · 시놉시스 보조)",
             "② 템플릿·출처 원칙을 '먼저' 정리",
             "③ 계산은 분리하고 '검증'",
             "④ 시놉시스 게이트를 교육의 중심에",
             "⑤ '초안 능력'보다 '검토 능력'을 가르쳐라",
         ], size=19),
         notes="우리 조직은 내일 무엇부터 할까요. 다섯 가지입니다. 첫째, 전체 자동화 말고 작은 과업부터 — 배경조사 요약, 비교표, 시놉시스 보조 같은 저위험 과업. 둘째, AI를 들이기 전에 사람이 쓸 기준부터 — 템플릿, 출처 원칙. 셋째, 계산은 글쓰기에서 분리하고 반드시 검증. 넷째, 시놉시스 승인 게이트를 교육 중심에. 다섯째, 가장 중요 — 누가 빨리 쓰는가보다 누가 잘 검토하는가를 가르치십시오. AI 시대의 핵심 역량입니다."),
    dict(sec="마무리", title="AI는 신입사원입니다. 잘 교육하면 동료가 됩니다.", accent=False,
         render=big("AI는 신입사원입니다.\n잘 교육하면 동료가 됩니다.",
                    "과장하지 말 것 · 과소평가하지 말 것 · 책임은 사람에게"),
         notes="처음 질문으로 돌아가겠습니다. 신입에게 첫날 계획서를 안 맡기듯, AI에게도 맥락과 구조를 갖춰줘야 합니다. AI는 신입사원입니다. 잘 교육하면 좋은 동료가 됩니다. 결국 임상개발에서 AI의 성패는 얼마나 유창한가가 아니라 얼마나 잘 구조화·검증되었는가에 달려 있습니다. 과장하지도 과소평가하지도 마십시오. 책임은 언제나 사람에게 있습니다. 좋은 AI는 우연히 생기지 않습니다 — 교육되고 구조화되고 검토되는 환경 속에서 만들어집니다. 감사합니다."),
]

for i, sd in enumerate(SLIDES, start=1):
    s = base(i, sd["sec"], sd["title"], accent=sd.get("accent", False))
    sd["render"](s)
    notes(s, sd["notes"])

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "presentation_genAI_protocol_DDI_ko.pptx")
prs.save(OUT)
print(f"saved: {OUT}  ({len(SLIDES)} slides)")
